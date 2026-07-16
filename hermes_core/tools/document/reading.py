# Copyright 2026 Kabuqina Contributors
# SPDX-License-Identifier: Apache-2.0

"""Document reading (Docling + format fallbacks), split from document_tools.py."""

from __future__ import annotations

import hashlib
import html
import io
import json
import logging
import os
import sys
import tempfile
import textwrap
import threading
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from kabuqina_constants import get_kabuqina_home
from tools.registry import registry, tool_error
from tools.document.common import (
    DocumentSpecError, Rgb, _PDF_BLOCK_TYPES, _PDF_TEMPLATES,
    _PPTX_SLIDE_LAYOUTS, _PPTX_SLIDE_TYPES, _desktop_workspace_root,
    _is_outside_workspace, _json, _power_user_reads_anywhere, _validate_read_path,
    _text, _list, _string_list, _dict,
)

logger = logging.getLogger(__name__)


_DOCLING_CONVERTERS: Dict[str, Any] = {}


_DOCLING_CONVERTER_LOCK = threading.RLock()


_DOCLING_IO_POOL: Optional["ThreadPoolExecutor"] = None


_DOCLING_IO_POOL_LOCK = threading.Lock()


_TORCH_PRIME_LOCK = threading.Lock()


def _get_docling_io_pool() -> "ThreadPoolExecutor":
    """Single-worker pool: all Docling/torch init and convert on one thread."""
    global _DOCLING_IO_POOL
    with _DOCLING_IO_POOL_LOCK:
        if _DOCLING_IO_POOL is None:
            from concurrent.futures import ThreadPoolExecutor

            _DOCLING_IO_POOL = ThreadPoolExecutor(max_workers=1, thread_name_prefix="docling-io")
        return _DOCLING_IO_POOL


def _run_on_docling_thread(fn, /, *args, **kwargs):
    return _get_docling_io_pool().submit(lambda: fn(*args, **kwargs)).result()


_DOCLING_SUFFIXES = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".html",
    ".htm",
    ".md",
    ".markdown",
    ".adoc",
    ".asciidoc",
    ".csv",
    ".png",
    ".jpg",
    ".jpeg",
    ".tif",
    ".tiff",
    ".bmp",
    ".webp",
    ".xml",
    ".wav",
    ".mp3",
    ".vtt",
}


_FALLBACK_SUFFIXES = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".html",
    ".htm",
    ".md",
    ".markdown",
    ".csv",
    ".txt",
}


_LIGHTWEIGHT_FIRST_SUFFIXES = {
    ".docx",
    ".pptx",
    ".xlsx",
    ".html",
    ".htm",
    ".md",
    ".markdown",
    ".csv",
    ".txt",
}


_KIND_BY_SUFFIX = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    ".html": "html",
    ".htm": "html",
    ".md": "markdown",
    ".markdown": "markdown",
    ".csv": "csv",
    ".txt": "text",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".tif": "image",
    ".tiff": "image",
    ".bmp": "image",
    ".webp": "image",
}


def _error_payload(message: str, **extra) -> Dict[str, Any]:
    return json.loads(tool_error(message, **extra))


def _document_kind(path: Path) -> str:
    return _KIND_BY_SUFFIX.get(path.suffix.lower(), path.suffix.lower().lstrip(".") or "unknown")


def _read_cache_root() -> Path:
    data_dir = (
        os.environ["KABUQINA_DATA_DIR"]
        if "KABUQINA_DATA_DIR" in os.environ
        else os.environ.get("HERMESDESK_DATA_DIR", "")
    )
    raw = (
        data_dir
        or str(get_kabuqina_home())
        or os.environ.get("LOCALAPPDATA")
        or tempfile.gettempdir()
    )
    root = Path(raw).expanduser()
    if root.name != "read-cache":
        root = root / "read-cache"
    root.mkdir(parents=True, exist_ok=True)
    return root


def _read_id_for_payload(payload: Dict[str, Any]) -> str:
    content = str(payload.get("content") or "")
    path = str(payload.get("path") or "")
    engine = str(payload.get("engine") or "")
    digest = hashlib.sha256()
    digest.update(path.encode("utf-8", errors="replace"))
    digest.update(b"\0")
    digest.update(engine.encode("utf-8", errors="replace"))
    digest.update(b"\0")
    digest.update(content.encode("utf-8", errors="replace"))
    return digest.hexdigest()[:24]


def _persist_read_result(payload: Dict[str, Any]) -> Tuple[str, str]:
    """Persist a Read-layer payload so downstream tools can consume it by reference."""
    read_id = str(payload.get("read_id") or _read_id_for_payload(payload))
    cache_path = _read_cache_root() / f"{read_id}.json"
    stored = dict(payload)
    stored["read_id"] = read_id
    stored["cache_path"] = str(cache_path)
    cache_path.write_text(_json(stored), encoding="utf-8")
    return read_id, str(cache_path)


def _load_read_result(read_id: str) -> Dict[str, Any]:
    safe_id = "".join(ch for ch in str(read_id or "") if ch.isalnum() or ch in {"-", "_"})
    if not safe_id or safe_id != str(read_id):
        raise ValueError("invalid read_id")
    cache_path = _read_cache_root() / f"{safe_id}.json"
    if not cache_path.exists():
        raise FileNotFoundError(f"read cache entry not found: {safe_id}")
    return json.loads(cache_path.read_text(encoding="utf-8"))


def _with_read_metadata(payload: Dict[str, Any], document_path: Path) -> Dict[str, Any]:
    metadata = payload.get("metadata") if isinstance(payload.get("metadata"), dict) else {}
    suffix = document_path.suffix.lower()
    enriched = {
        **payload,
        "path": str(document_path),
        "metadata": {
            **metadata,
            "kind": metadata.get("kind") or _document_kind(document_path),
            "suffix": suffix,
            "source_name": document_path.name,
            "parser_engine": payload.get("engine", ""),
        },
    }
    return enriched


def _finalize_read_payload(payload: Dict[str, Any], document_path: Path, *, include_content: bool) -> Dict[str, Any]:
    enriched = _with_read_metadata(payload, document_path)
    content = str(enriched.get("content") or "")
    read_id, cache_path = _persist_read_result(enriched)
    content_value = enriched.pop("content", "")
    enriched["read_id"] = read_id
    enriched["cache_path"] = cache_path
    enriched["content_chars"] = len(content)
    enriched["metadata"]["read_id"] = read_id
    enriched["content_hint"] = (
        "Full extracted content is stored in read-cache. Use read_file on cache_path "
        "with offset/limit to inspect the cached JSON content, or pass read_id to "
        "material_index_build. For extracting formulas from PDFs, use pdf_read_precise "
        "or document_read_precise with mode=math and read the extracted Markdown; do "
        "not use vision_analyze unless the PDF reader failed or the user explicitly "
        "asked for visual image description."
    )
    if not include_content:
        enriched["content"] = ""
        enriched["content_omitted"] = True
    else:
        enriched["content"] = content_value
    return enriched


def _resolve_docling_artifacts_path(profile: str = "fast") -> Optional[Path]:
    """Return bundled/offline Docling model dir when present."""
    try:
        from docling_math_models import resolve_docling_artifacts_path

        resolved = resolve_docling_artifacts_path(profile=profile)
        if resolved is not None:
            return resolved
        if profile == "math":
            return None
    except ImportError:
        pass

    explicit = os.environ.get("DOCLING_ARTIFACTS_PATH", "").strip()
    if explicit:
        path = Path(explicit).expanduser()
        if path.is_dir():
            return path

    bundle_dir = os.environ.get("HERMESDESK_BUNDLE_DIR", "").strip()
    if bundle_dir:
        path = Path(bundle_dir) / "docling-models"
        if path.is_dir():
            return path
    return None


def _code_formula_bundle_present(artifacts_path: Optional[Path]) -> bool:
    if artifacts_path is None:
        return False
    formula_dir = artifacts_path / "ds4sd--CodeFormula"
    if not formula_dir.is_dir():
        return False
    return any(formula_dir.rglob("*.safetensors")) or any(formula_dir.rglob("*.bin"))


def _require_math_artifacts_bundled() -> None:
    """Non-desktop fallback when docling_math_models is unavailable."""
    artifacts_path = _resolve_docling_artifacts_path("fast")
    if not _code_formula_bundle_present(artifacts_path):
        target = (
            str(artifacts_path / "ds4sd--CodeFormula")
            if artifacts_path is not None
            else "<bundle>/docling-models/ds4sd--CodeFormula"
        )
        raise ValueError(
            "mode=math requires offline CodeFormula weights at "
            f"{target}. Re-run .\\python\\build_bundle.ps1 to bundle ds4sd/CodeFormula, "
            "or set DOCLING_ARTIFACTS_PATH to a directory that contains ds4sd--CodeFormula/."
        )


def _ensure_math_artifacts() -> None:
    """Ensure CodeFormula weights exist; fail with settings hint when missing."""
    try:
        from docling_math_models import ensure_code_formula_available_for_math

        ensure_code_formula_available_for_math()
        return
    except ImportError:
        pass
    _require_math_artifacts_bundled()


def _docling_profile_for_mode(mode: str) -> str:
    key = (mode or "").strip().lower()
    if key in {"precise", "math"}:
        return key
    return "fast"


def _docling_ocr_enabled() -> bool:
    """OCR is off by default — it loads an extra CPU model and runs a per-page
    pass that born-digital PDFs (papers, task books) never need, and it is the
    single biggest slowdown of precise/math reads on machines without a GPU. Set
    HERMESDESK_DOCLING_OCR=1 to re-enable it for genuinely scanned documents."""
    return _text(os.getenv("HERMESDESK_DOCLING_OCR")).lower() in {"1", "true", "yes", "on"}


def _docling_max_pages() -> Optional[int]:
    """Optional hard cap on pages handed to Docling, so a huge PDF cannot pin the
    CPU for an unbounded time. Unset by default (no cap); set
    HERMESDESK_DOCLING_MAX_PAGES=N to bound it."""
    raw = _text(os.getenv("HERMESDESK_DOCLING_MAX_PAGES"))
    if not raw:
        return None
    try:
        value = int(raw)
    except ValueError:
        return None
    return value if value > 0 else None


def _docling_math_max_pages() -> Optional[int]:
    """Default guard for CPU-bound CodeFormula extraction.

    ``mode=math`` runs CodeFormula per page and is very slow without a GPU. Keep
    the default small so agents must request a focused page range instead of
    accidentally running an entire PDF. Set HERMESDESK_DOCLING_MATH_MAX_PAGES=0
    to allow full-document math extraction.
    """
    raw = _text(os.getenv("HERMESDESK_DOCLING_MATH_MAX_PAGES", "2")).lower()
    if raw in {"0", "false", "no", "off", "none", "unlimited"}:
        return None
    try:
        value = int(raw)
    except ValueError:
        return 2
    return value if value > 0 else None


def _docling_tables_enabled() -> bool:
    """Table-structure (TableFormer) is the heaviest per-table CPU stage. It stays
    on for mode=precise but is off for mode=math, whose job is formula extraction,
    not table reconstruction. Set HERMESDESK_DOCLING_TABLES=1 to force it on for
    math too."""
    return _text(os.getenv("HERMESDESK_DOCLING_TABLES")).lower() in {"1", "true", "yes", "on"}


def _prime_torch_for_docling() -> None:
    with _TORCH_PRIME_LOCK:
        import torch
        import torch.library  # noqa: F401
        if not hasattr(torch, "library"):
            raise AttributeError("torch imported without torch.library")


def _configure_pdf_pipeline_options(pipeline_options: Any, profile: str) -> None:
    if profile == "fast":
        pipeline_options.do_ocr = False
        pipeline_options.do_table_structure = False
        if hasattr(pipeline_options, "force_backend_text"):
            pipeline_options.force_backend_text = True
        if hasattr(pipeline_options, "do_formula_enrichment"):
            pipeline_options.do_formula_enrichment = False
        if hasattr(pipeline_options, "do_code_enrichment"):
            pipeline_options.do_code_enrichment = False
        return

    if profile == "precise":
        pipeline_options.do_ocr = _docling_ocr_enabled()
        if hasattr(pipeline_options, "do_formula_enrichment"):
            pipeline_options.do_formula_enrichment = False
        if hasattr(pipeline_options, "do_code_enrichment"):
            pipeline_options.do_code_enrichment = False
        return

    if profile == "math":
        pipeline_options.do_ocr = _docling_ocr_enabled()
        # math is the *formula* profile: skip TableFormer (the heaviest per-table
        # CPU stage) so formula extraction is not throttled by table-structure
        # inference the caller did not ask for. Use mode=precise for tables, or
        # set HERMESDESK_DOCLING_TABLES=1 to keep tables in math too.
        if hasattr(pipeline_options, "do_table_structure"):
            pipeline_options.do_table_structure = _docling_tables_enabled()
        if not hasattr(pipeline_options, "do_formula_enrichment"):
            raise ValueError("Docling PdfPipelineOptions does not support formula enrichment")
        pipeline_options.do_formula_enrichment = True
        if hasattr(pipeline_options, "do_code_enrichment"):
            pipeline_options.do_code_enrichment = False
        if hasattr(pipeline_options, "code_formula_options"):
            pipeline_options.code_formula_options.extract_formulas = True
            pipeline_options.code_formula_options.extract_code = False
        return

    raise ValueError(f"unknown Docling profile: {profile}")


def _create_docling_converter(profile: str = "fast"):
    if profile == "math":
        _ensure_math_artifacts()
    _prime_torch_for_docling()

    from docling.datamodel.base_models import InputFormat  # type: ignore
    from docling.datamodel.pipeline_options import PdfPipelineOptions  # type: ignore
    from docling.document_converter import DocumentConverter, PdfFormatOption  # type: ignore

    pipeline_options = PdfPipelineOptions()
    artifacts_path = _resolve_docling_artifacts_path(profile)
    if artifacts_path is not None:
        pipeline_options.artifacts_path = artifacts_path
        easyocr_dir = artifacts_path / "EasyOcr"
        if easyocr_dir.is_dir() and hasattr(pipeline_options, "ocr_options"):
            pipeline_options.ocr_options.model_storage_directory = str(easyocr_dir)
            pipeline_options.ocr_options.download_enabled = False

    _configure_pdf_pipeline_options(pipeline_options, profile)

    return DocumentConverter(
        format_options={
            InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options),
        }
    )


def _get_docling_converter(profile: str = "fast"):
    profile = _docling_profile_for_mode(profile)
    if profile not in _DOCLING_CONVERTERS:
        with _DOCLING_CONVERTER_LOCK:
            if profile not in _DOCLING_CONVERTERS:
                _DOCLING_CONVERTERS[profile] = _create_docling_converter(profile)
    return _DOCLING_CONVERTERS[profile]


def reset_docling_converter_cache() -> None:
    with _DOCLING_CONVERTER_LOCK:
        _DOCLING_CONVERTERS.clear()


def warm_docling_converter(mode: str = "auto") -> Dict[str, Any]:
    """Eagerly initialize Docling so the first user document read is faster."""
    return _run_on_docling_thread(_warm_docling_converter_impl, mode)


def _warm_docling_converter_impl(mode: str = "auto") -> Dict[str, Any]:
    profile = _docling_profile_for_mode(mode)
    converter = _get_docling_converter(profile)
    return {
        "ok": True,
        "engine": "docling",
        "profile": profile,
        "converter": type(converter).__name__,
        "artifacts_path": str(_resolve_docling_artifacts_path() or ""),
    }


def _build_docling_converter():
    """Compatibility wrapper for older call sites/tests."""
    return _get_docling_converter("fast")


def _is_unsupported_runtime_error(lowered: str) -> bool:
    """Detect Docling failures caused by an unsupported Python/ML runtime.

    Symptoms seen on unsupported interpreters (e.g. Python 3.14 with a torch
    build that has not caught up): torch import is "partially initialized" /
    missing attributes, or the pdfium backend cannot agree on a page count.
    These are environment problems, not bad input, so the hint must point at
    the supported runtime rather than at the document.
    """
    torch_broken = "torch" in lowered and any(
        token in lowered
        for token in ("partially initialized", "has no attribute", "no attribute 'library'")
    )
    pdfium_broken = "inconsistent number of pages" in lowered or "pdfium" in lowered
    torch_dispatcher_broken = (
        "already a kernel registered" in lowered
        and ("wait_tensor" in lowered or "_c10d_functional" in lowered)
    )
    return torch_broken or pdfium_broken or torch_dispatcher_broken


def _is_torch_dispatcher_registration_error(lowered: str) -> bool:
    return (
        "already a kernel registered" in lowered
        and ("wait_tensor" in lowered or "_c10d_functional" in lowered)
    )


def _format_docling_error(exc: BaseException) -> str:
    msg = str(exc).strip() or type(exc).__name__
    lowered = msg.lower()
    if _is_unsupported_runtime_error(lowered):
        runtime = f"CPython {sys.version_info.major}.{sys.version_info.minor} ({sys.executable})"
        if _is_torch_dispatcher_registration_error(lowered):
            hint = (
                "This points to duplicate PyTorch dispatcher registration in the current "
                "process, usually after a failed Docling/PyTorch warmup or re-import. Fully "
                "restart Kabuqina once; the background Docling warmup is disabled by default "
                "to avoid poisoning the main agent process."
            )
        elif sys.version_info[:2] == (3, 11):
            hint = (
                "The bundled 3.11 runtime is active, so this is most likely torch/Docling "
                "initialization state left behind in the current process. Fully restart "
                "Kabuqina once; if it repeats, keep the text fallback and report this "
                "docling_error."
            )
        else:
            hint = (
                "Kabuqina's supported runtime is the bundled CPython 3.11 from "
                "python/build_bundle.ps1; torch/pdfium may fail on newer interpreters "
                "such as 3.14. Run on the bundled 3.11 runtime to restore precise reads."
            )
        return (
            f"Docling backend failed on this Python/ML runtime ({type(exc).__name__}): {msg}. "
            f"This is an environment problem, not a bad document. Runtime: {runtime}. "
            f"{hint} A text-only fallback was used so this read could still succeed."
        )
    if any(token in lowered for token in ("huggingface", "timed out", "timeout", "connection")):
        return (
            f"Docling model load failed ({type(exc).__name__}): {msg}. "
            "Offline models may be missing from the app bundle, or HuggingFace is unreachable."
        )
    if isinstance(exc, ImportError):
        return f"Docling is not installed: {msg}"
    try:
        from docling_math_models import CodeFormulaMissingError

        if isinstance(exc, CodeFormulaMissingError):
            return msg
    except ImportError:
        pass
    if "code_formula_model_missing" in lowered:
        return msg
    if isinstance(exc, PermissionError) or "declined the formula model download" in lowered:
        return (
            f"Docling formula model unavailable ({type(exc).__name__}): {msg}. "
            "Open Kabuqina Settings → Load packages to download (~500 MB), then retry mode=math."
        )
    if "codeformula" in lowered and ("mode=math requires" in lowered or "hfvalidationerror" in lowered):
        return (
            f"Docling formula model unavailable ({type(exc).__name__}): {msg}. "
            "Open Kabuqina Settings → Load packages to download ds4sd/CodeFormula (~500 MB), "
            "then retry mode=math."
        )
    return f"Docling failed ({type(exc).__name__}): {msg}"


def _read_with_docling(
    document_path: Path,
    mode: str,
    page_range: Optional[Tuple[int, int]] = None,
) -> Dict[str, Any]:
    return _run_on_docling_thread(_read_with_docling_impl, document_path, mode, page_range)


def _read_with_docling_page_range(
    document_path: Path,
    mode: str,
    page_range: Optional[Tuple[int, int]],
) -> Dict[str, Any]:
    if page_range is None:
        return _read_with_docling(document_path, mode)
    return _read_with_docling(document_path, mode, page_range)


def _read_with_docling_impl(
    document_path: Path,
    mode: str,
    page_range: Optional[Tuple[int, int]] = None,
) -> Dict[str, Any]:
    profile = _docling_profile_for_mode(mode)
    converter = _get_docling_converter(profile)
    convert_kwargs: Dict[str, Any] = {}
    max_pages = _docling_max_pages()
    if max_pages is not None:
        convert_kwargs["max_num_pages"] = max_pages
    if page_range is not None:
        convert_kwargs["page_range"] = page_range
    result = converter.convert(str(document_path), **convert_kwargs)
    document = result.document
    markdown = document.export_to_markdown()
    return {
        "ok": True,
        "engine": "docling",
        "mode": mode,
        "profile": profile,
        "path": str(document_path),
        "pages": len(getattr(document, "pages", []) or []),
        "content": markdown,
    }


def _read_pdf_with_pypdf(pdf_path: Path) -> Dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(str(pdf_path))
    chunks: List[str] = []
    for idx, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        chunks.append(f"<!-- page:{idx} -->\n{text}".strip())
    return {
        "ok": True,
        "engine": "pypdf",
        "mode": "fallback",
        "path": str(pdf_path),
        "pages": len(reader.pages),
        "content": "\n\n".join(chunks).strip(),
        "warning": "Used text-only pypdf fallback because Docling could not parse this PDF.",
    }


def _pdf_page_count(pdf_path: Path) -> Optional[int]:
    try:
        from pypdf import PdfReader

        return len(PdfReader(str(pdf_path)).pages)
    except Exception:
        return None


def _read_docx_with_python_docx(docx_path: Path) -> Dict[str, Any]:
    from docx import Document

    doc = Document(str(docx_path))
    chunks: List[str] = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                chunks.append(" | ".join(cells))
    return {
        "ok": True,
        "engine": "python-docx",
        "mode": "fallback",
        "path": str(docx_path),
        "pages": 0,
        "content": "\n\n".join(chunks).strip(),
        "warning": "Used text-only python-docx fallback because Docling could not parse this document.",
    }


def _read_pptx_with_python_pptx(pptx_path: Path) -> Dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slides: List[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        parts: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        slides.append(f"<!-- slide:{idx} -->\n" + "\n\n".join(parts).strip())
    return {
        "ok": True,
        "engine": "python-pptx",
        "mode": "fallback",
        "path": str(pptx_path),
        "slides": len(prs.slides),
        "content": "\n\n".join(slides).strip(),
        "warning": "Used text-only python-pptx fallback because Docling could not parse this document.",
    }


def _read_xlsx_with_openpyxl(xlsx_path: Path) -> Dict[str, Any]:
    from openpyxl import load_workbook

    wb = load_workbook(str(xlsx_path), read_only=True, data_only=True)
    chunks: List[str] = []
    for ws in wb.worksheets:
        chunks.append(f"## {ws.title}")
        for row in ws.iter_rows(values_only=True):
            values = ["" if cell is None else str(cell) for cell in row]
            if any(value.strip() for value in values):
                chunks.append(" | ".join(values))
    return {
        "ok": True,
        "engine": "openpyxl",
        "mode": "fallback",
        "path": str(xlsx_path),
        "sheets": len(wb.worksheets),
        "content": "\n".join(chunks).strip(),
        "warning": "Used text-only openpyxl fallback because Docling could not parse this document.",
    }


def _read_plain_text(document_path: Path) -> Dict[str, Any]:
    content = document_path.read_text(encoding="utf-8", errors="replace")
    return {
        "ok": True,
        "engine": "text",
        "mode": "fallback",
        "path": str(document_path),
        "content": content,
        "warning": "Used plain-text fallback because Docling could not parse this document.",
    }


def _read_with_fallback(document_path: Path) -> Optional[Dict[str, Any]]:
    suffix = document_path.suffix.lower()
    if suffix == ".pdf":
        return _read_pdf_with_pypdf(document_path)
    if suffix == ".docx":
        return _read_docx_with_python_docx(document_path)
    if suffix == ".pptx":
        return _read_pptx_with_python_pptx(document_path)
    if suffix == ".xlsx":
        return _read_xlsx_with_openpyxl(document_path)
    if suffix in {".html", ".htm", ".md", ".markdown", ".csv", ".txt"}:
        return _read_plain_text(document_path)
    return None


def _should_try_lightweight_first(suffix: str, mode: str) -> bool:
    return _docling_profile_for_mode(mode) == "fast" and suffix in _LIGHTWEIGHT_FIRST_SUFFIXES


def _pdf_fast_text_is_sufficient(payload: Dict[str, Any]) -> bool:
    """Is the pypdf text good enough to skip Docling? Guards against scanned PDFs.

    A text-based PDF (most papers/reports) extracts plenty of characters; a
    scanned/image PDF yields almost nothing, in which case we fall through to
    Docling so layout/OCR still has a chance.
    """
    content = str(payload.get("content") or "").strip()
    if len(content) < 200:
        return False
    pages = int(payload.get("pages") or 0)
    if pages > 0 and (len(content) / pages) < 40:
        return False
    return True


def _attach_docling_error(payload: Dict[str, Any], docling_error: Optional[str]) -> Dict[str, Any]:
    if docling_error:
        payload["docling_error"] = docling_error
    return payload


def _should_avoid_docling_fallback(mode: str, exc: Optional[BaseException] = None) -> bool:
    if _docling_profile_for_mode(mode) == "math":
        return True
    if exc is not None:
        lowered = str(exc).lower()
        if "code_formula_model_missing" in lowered:
            return True
    try:
        from docling_math_models import CodeFormulaMissingError

        if isinstance(exc, CodeFormulaMissingError):
            return True
    except ImportError:
        pass
    if isinstance(exc, PermissionError):
        return True
    if exc is not None and "declined the formula model download" in str(exc).lower():
        return True
    return False


def _normalize_page_range(
    page_start: Optional[int] = None,
    page_end: Optional[int] = None,
) -> Tuple[Optional[Tuple[int, int]], Optional[Dict[str, Any]]]:
    if page_start is None and page_end is None:
        return None, None
    try:
        start = int(page_start) if page_start is not None else 1
        end = int(page_end) if page_end is not None else start
    except (TypeError, ValueError):
        return None, _error_payload(
            "page_start/page_end must be positive integers.",
            ok=False,
            code="invalid_page_range",
        )
    if start < 1 or end < 1 or end < start:
        return None, _error_payload(
            "Invalid page range: use 1-based page_start/page_end with page_end >= page_start.",
            ok=False,
            code="invalid_page_range",
        )
    return (start, end), None


def _guard_math_page_count(
    document_path: Path,
    mode: str,
    page_range: Optional[Tuple[int, int]],
) -> Optional[Dict[str, Any]]:
    if document_path.suffix.lower() != ".pdf" or _docling_profile_for_mode(mode) != "math":
        return None
    max_pages = _docling_math_max_pages()
    if max_pages is None:
        return None
    if page_range is not None:
        selected_pages = page_range[1] - page_range[0] + 1
        if selected_pages <= max_pages:
            return None
        return _error_payload(
            (
                f"mode=math selected {selected_pages} pages, above the current "
                f"CodeFormula CPU guard of {max_pages} pages."
            ),
            ok=False,
            code="docling_math_too_many_pages",
            pages=selected_pages,
            max_pages=max_pages,
            hint=(
                "Use a smaller page_start/page_end range, split the extraction into "
                "several calls, or set HERMESDESK_DOCLING_MATH_MAX_PAGES=0 to allow "
                "full-document math extraction."
            ),
        )
    page_count = _pdf_page_count(document_path)
    if page_count is None or page_count <= max_pages:
        return None
    return _error_payload(
        (
            f"mode=math would run CodeFormula over {page_count} PDF pages on CPU. "
            f"The default guard allows {max_pages} pages per call."
        ),
        ok=False,
        code="docling_math_too_many_pages",
        pages=page_count,
        max_pages=max_pages,
        hint=(
            "Pass page_start/page_end for the pages that contain formulas, split the "
            "document into smaller ranges, or set HERMESDESK_DOCLING_MATH_MAX_PAGES=0 "
            "to allow full-document math extraction."
        ),
    )


def _read_document_precise_payload(
    document_path: Path,
    *,
    mode: str,
    include_content: bool = True,
    page_start: Optional[int] = None,
    page_end: Optional[int] = None,
) -> Dict[str, Any]:
    suffix = document_path.suffix.lower()
    docling_error: Optional[str] = None
    page_range, page_error = _normalize_page_range(page_start, page_end)
    if page_error is not None:
        return page_error

    if suffix == ".doc":
        return _error_payload(
            "Legacy .doc files are not directly supported by Docling.",
            ok=False,
            code="unsupported_legacy_doc",
            hint="Convert the file with LibreOffice headless to .docx or .pdf, then read it again.",
        )
    if suffix not in _DOCLING_SUFFIXES and suffix not in _FALLBACK_SUFFIXES:
        return _error_payload(
            f"Unsupported document type: {suffix or '(no extension)'}",
            ok=False,
            code="unsupported_document_type",
            supported=sorted(_DOCLING_SUFFIXES | _FALLBACK_SUFFIXES),
        )

    # Ensure the optional CodeFormula pack is available *before* the page-count
    # guard. The guard only caps CPU-bound formula inference; it must not gate the
    # one-time, approval-prompted model download. So attempting math extraction at
    # all should offer the download — even for a large PDF the guard will then ask
    # to narrow to a page range. Once the pack is present the page-ranged retry
    # runs without re-downloading.
    if suffix in _DOCLING_SUFFIXES and _docling_profile_for_mode(mode) == "math":
        try:
            _ensure_math_artifacts()
        except Exception as exc:
            docling_error = _format_docling_error(exc)
            return _error_payload(
                docling_error or "Docling formula model unavailable.",
                ok=False,
                code="docling_math_unavailable",
                docling_error=docling_error,
            )

    math_guard = _guard_math_page_count(document_path, mode, page_range)
    if math_guard is not None:
        return math_guard

    # Fast text-first for PDFs in auto/fast mode: a text-based paper/report
    # extracts in well under a second with pypdf, versus minutes of Docling
    # layout inference on CPU. Scanned/sparse PDFs fail the sufficiency guard
    # and fall through to Docling. precise/math modes always use Docling.
    if suffix == ".pdf" and _docling_profile_for_mode(mode) == "fast":
        try:
            fast = _read_pdf_with_pypdf(document_path)
        except Exception:
            fast = None
        if fast is not None and _pdf_fast_text_is_sufficient(fast):
            fast["mode"] = mode
            fast["engine"] = "pypdf"
            fast["warning"] = (
                "Fast text-only PDF read (pypdf) — chosen for speed, not a Docling "
                "failure. For layout, tables, or formula extraction, re-read with "
                "mode=precise or mode=math."
            )
            return _finalize_read_payload(fast, document_path, include_content=include_content)

    if _should_try_lightweight_first(suffix, mode):
        try:
            fallback = _read_with_fallback(document_path)
            if fallback is not None:
                fallback["mode"] = mode
                return _finalize_read_payload(fallback, document_path, include_content=include_content)
        except Exception:
            pass

    if suffix in _DOCLING_SUFFIXES:
        try:
            return _finalize_read_payload(
                _read_with_docling_page_range(document_path, mode, page_range),
                document_path,
                include_content=include_content,
            )
        except Exception as exc:
            docling_error = _format_docling_error(exc)
            if _should_avoid_docling_fallback(mode, exc):
                return _error_payload(
                    docling_error or "Docling read failed.",
                    ok=False,
                    code="docling_math_unavailable",
                    docling_error=docling_error,
                )

    try:
        fallback = _read_with_fallback(document_path)
        if fallback is not None:
            return _finalize_read_payload(
                _attach_docling_error(fallback, docling_error),
                document_path,
                include_content=include_content,
            )
    except Exception as exc:
        return _error_payload(
            f"Document fallback read failed: {exc}",
            ok=False,
            docling_error=docling_error,
        )

    return _error_payload(
        "Docling could not parse this document and no local fallback is available.",
        ok=False,
        code="docling_failed_no_fallback",
        docling_error=docling_error,
    )


def pdf_read_precise(
    path: str,
    mode: str = "auto",
    include_content: bool = True,
    page_start: Optional[int] = None,
    page_end: Optional[int] = None,
) -> str:
    pdf_path = Path(path).expanduser()
    validation_error = _validate_read_path(pdf_path, path, "PDF")
    if validation_error is not None:
        return validation_error
    if pdf_path.suffix.lower() != ".pdf":
        return tool_error("pdf_read_precise only accepts .pdf files")
    read_kwargs: Dict[str, Any] = {
        "mode": mode,
        "include_content": include_content,
    }
    if page_start is not None:
        read_kwargs["page_start"] = page_start
    if page_end is not None:
        read_kwargs["page_end"] = page_end
    return _json(
        _read_document_precise_payload(
            pdf_path,
            **read_kwargs,
        )
    )


def document_read_precise(
    path: str,
    mode: str = "auto",
    include_content: bool = True,
    page_start: Optional[int] = None,
    page_end: Optional[int] = None,
) -> str:
    document_path = Path(path).expanduser()
    validation_error = _validate_read_path(document_path, path, "Document")
    if validation_error is not None:
        return validation_error
    read_kwargs: Dict[str, Any] = {
        "mode": mode,
        "include_content": include_content,
    }
    if page_start is not None:
        read_kwargs["page_start"] = page_start
    if page_end is not None:
        read_kwargs["page_end"] = page_end
    return _json(
        _read_document_precise_payload(
            document_path,
            **read_kwargs,
        )
    )
